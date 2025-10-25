import os
import json
import logging
from typing import Dict, List, Any, Tuple
from dataclasses import dataclass
import anthropic
from pathlib import Path
import pypdf
from docx import Document
import pptx
import asyncio
from concurrent.futures import ThreadPoolExecutor
import re

logger = logging.getLogger(__name__)

@dataclass
class ExtractedKnowledge:
    company_info: Dict[str, Any]
    sales_approach: str
    products: List[Dict[str, Any]]
    key_messages: List[str]
    value_propositions: List[str]
    target_audience: Dict[str, Any]
    competitive_advantages: List[str]

class KnowledgeExtractionAgent:
    """
    Agent that uses Claude AI to extract structured knowledge from company documents.
    Specializes in understanding company voice, sales approach, and product information.
    """
    
    def __init__(self):
        self.claude_client = anthropic.Anthropic(
            api_key=os.getenv('ANTHROPIC_API_KEY')
        )
        
    def extract_knowledge_from_files(self, file_paths: List[str], document_type: str = None) -> Dict[str, Any]:
        """
        Extract structured knowledge from uploaded documents using Claude AI.
        Now supports chunking for large documents and parallel processing.
        
        Args:
            file_paths: List of file paths to process
            document_type: Optional document type to guide extraction
        """
        try:
            logger.info(f"Starting knowledge extraction from {len(file_paths)} files")
            if document_type:
                logger.info(f"Using explicit document type: {document_type}")
            
            # Check if we should use parallel processing
            use_parallel = len(file_paths) > 1
            
            if use_parallel:
                logger.info("Using parallel processing for multiple files")
                return self.extract_knowledge_parallel(file_paths, document_type)
            
            # Single file processing with chunking support
            documents_content = []
            for file_path in file_paths:
                content = self._read_document(file_path)
                if content:
                    # Check if document needs chunking
                    if len(content) > 100000:  # ~100k characters
                        logger.info(f"Large document detected ({len(content)} chars), applying chunking")
                        chunks = self._chunk_large_document(content, file_path)
                        documents_content.extend(chunks)
                    else:
                        documents_content.append({
                            'filename': Path(file_path).name,
                            'content': content,
                            'is_chunk': False,
                            'chunk_index': 0
                        })
            
            if not documents_content:
                return {
                    "success": False,
                    "error": "No readable documents found"
                }
            
            # Use Claude to extract knowledge with optional document type
            extracted_knowledge = self._extract_with_claude(documents_content, document_type)
            
            # Add quality validation and confidence scoring
            validated_knowledge = self._validate_and_score_knowledge(extracted_knowledge, documents_content, document_type)
            
            logger.info("Knowledge extraction completed successfully")
            return {
                "success": True,
                "knowledge": validated_knowledge,
                "processing_method": "chunked" if any(doc.get('is_chunk', False) for doc in documents_content) else "standard",
                "quality_metrics": validated_knowledge.get("quality_metrics", {})
            }
            
        except Exception as e:
            logger.error(f"Knowledge extraction error: {e}")
            return {
                "success": False,
                "error": str(e)
            }
    
    def _read_document(self, file_path: str) -> str:
        """
        Read document content based on file type.
        """
        try:
            file_extension = Path(file_path).suffix.lower()
            
            if file_extension == '.txt':
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()
            elif file_extension == '.pdf':
                return self._read_pdf(file_path)
            elif file_extension in ['.doc', '.docx']:
                return self._read_word_document(file_path)
            elif file_extension in ['.ppt', '.pptx']:
                return self._read_powerpoint(file_path)
            else:
                logger.warning(f"Unsupported file type: {file_extension}")
                return None
                
        except Exception as e:
            logger.error(f"Error reading document {file_path}: {e}")
            return None
    
    def _read_pdf(self, file_path: str) -> str:
        """Read PDF content using pypdf."""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = pypdf.PdfReader(file)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
                return text.strip()
        except Exception as e:
            logger.error(f"Error reading PDF {file_path}: {e}")
            return f"[Error reading PDF: {str(e)}]"
    
    def _read_word_document(self, file_path: str) -> str:
        """Read Word document content using python-docx."""
        try:
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"Error reading Word document {file_path}: {e}")
            return f"[Error reading Word document: {str(e)}]"
    
    def _read_powerpoint(self, file_path: str) -> str:
        """Read PowerPoint content using python-pptx."""
        try:
            prs = pptx.Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"Error reading PowerPoint {file_path}: {e}")
            return f"[Error reading PowerPoint: {str(e)}]"
    
    def _extract_with_claude(self, documents: List[Dict[str, str]], document_type: str = None) -> Dict[str, Any]:
        """
        Use Claude AI to extract structured knowledge from documents.
        
        Args:
            documents: List of document content
            document_type: Optional explicit document type to guide extraction
        """
        try:
            # Prepare documents for Claude
            documents_text = "\n\n".join([
                f"Document: {doc['filename']}\n{doc['content']}"
                for doc in documents
            ])
            
            # Create the extraction prompt with optional document type
            prompt = self._build_extraction_prompt(documents_text, document_type)
            
            # Call Claude API
            response = self.claude_client.messages.create(
                model="claude-3-haiku-20240307",
                max_tokens=4000,
                messages=[{
                    "role": "user",
                    "content": prompt
                }]
            )
            
            # Parse the response
            extracted_text = response.content[0].text
            knowledge = self._parse_claude_response(extracted_text)
            
            return knowledge
            
        except Exception as e:
            logger.error(f"Claude extraction error: {e}")
            # Return default structure if Claude fails
            return self._get_default_knowledge_structure()
    
    def _build_extraction_prompt(self, documents_text: str, document_type: str = None) -> str:
        """
        Build the prompt for Claude to extract knowledge based on document type.
        
        Args:
            documents_text: The content of documents to analyze
            document_type: Optional explicit document type to guide extraction
        """
        if document_type:
            # Use explicit document type - skip detection step
            return self._build_targeted_prompt(documents_text, document_type)
        else:
            # Use automatic detection (existing logic)
            return self._build_detection_prompt(documents_text)
    
    def _build_targeted_prompt(self, documents_text: str, document_type: str) -> str:
        """
        Build a targeted prompt for a specific document type.
        """
        if document_type == "company_info":
            return f"""
You are an expert business analyst specializing in extracting company information from documents.

Extract the following information from the documents:

Documents to analyze:
{documents_text}

Extract company information in this JSON format:
{{
  "document_type": "company_info",
  "company_info": {{
    "company_name": "Name of the company",
    "industry": "Primary industry or sector",
    "company_size": "Size category (startup, SMB, enterprise)",
    "mission": "Company mission statement",
    "values": ["Core company values"],
    "founding_year": "Year founded if mentioned",
    "headquarters": "Location if mentioned"
  }},
  "sales_approach": "Detailed description of the company's sales methodology, approach, and philosophy",
  "products": [
    {{
      "name": "Product/service name",
      "description": "Detailed description",
      "target_customers": "Who this product serves",
      "key_features": ["Main features"],
      "benefits": ["Key benefits"],
      "pricing_model": "How it's priced if mentioned"
    }}
  ],
  "key_messages": [
    "Key messages the company uses in sales and marketing",
    "Value propositions",
    "Unique selling points"
  ],
  "value_propositions": [
    "Primary value propositions",
    "What makes them different",
    "Key benefits they provide"
  ],
  "target_audience": {{
    "primary_customers": "Description of ideal customers",
    "industries": ["Industries they serve"],
    "company_sizes": ["Target company sizes"],
    "pain_points": ["Problems they solve"]
  }},
  "competitive_advantages": [
    "What makes them unique",
    "Competitive differentiators",
    "Strengths vs competitors"
  ]
}}

IMPORTANT INSTRUCTIONS:
1. If information is not available, use "Not specified" or leave arrays empty
2. Focus on extracting actionable information that would help AI sales agents
3. Return only valid JSON, no additional text

Return the JSON structure.
"""
        
        elif document_type == "products":
            return f"""
You are an expert business analyst specializing in extracting product information from documents.

Extract the following information from the documents:

Documents to analyze:
{documents_text}

Extract product information in this JSON format:
{{
  "document_type": "products",
  "products": [
    {{
      "name": "Product/service name",
      "description": "Detailed description",
      "target_customers": "Who this product serves",
      "key_features": ["Main features"],
      "benefits": ["Key benefits"],
      "pricing_model": "How it's priced if mentioned",
      "use_cases": ["Specific use cases"],
      "integration_requirements": ["Technical requirements"],
      "competitors": ["Competing products/services"]
    }}
  ],
  "value_propositions": [
    "Primary value propositions",
    "What makes them different",
    "Key benefits they provide"
  ],
  "target_audience": {{
    "primary_customers": "Description of ideal customers",
    "industries": ["Industries they serve"],
    "company_sizes": ["Target company sizes"],
    "pain_points": ["Problems they solve"]
  }},
  "competitive_advantages": [
    "What makes them unique",
    "Competitive differentiators",
    "Strengths vs competitors"
  ],
  "key_messages": [
    "Key messages about the products",
    "Value propositions",
    "Unique selling points"
  ]
}}

IMPORTANT INSTRUCTIONS:
1. If information is not available, use "Not specified" or leave arrays empty
2. Focus on extracting actionable information that would help AI sales agents
3. Return only valid JSON, no additional text

Return the JSON structure.
"""
        
        elif document_type == "sales_training":
            return f"""
You are an expert business analyst specializing in extracting sales training knowledge from documents.

Extract the following information from the documents:

Documents to analyze:
{documents_text}

Extract sales training knowledge in this JSON format:
{{
  "document_type": "sales_training",
  "sales_methodologies": [
    "Sales techniques and approaches mentioned",
    "Best practices for prospecting",
    "Communication strategies"
  ],
  "target_audience_insights": {{
    "buyer_personas": ["Types of buyers mentioned"],
    "pain_points": ["Common problems prospects face"],
    "decision_making_process": ["How buyers make decisions"]
  }},
  "sales_philosophy": "Overall sales approach and mindset",
  "key_messages": [
    "Important sales messages and frameworks",
    "Value propositions that work",
    "Objection handling techniques"
  ],
  "industry_context": {{
    "industries_mentioned": ["Industries discussed"],
    "market_trends": ["Trends and insights"],
    "competitive_landscape": ["Competitive information"]
  }},
  "practical_advice": [
    "Actionable sales tips",
    "Scripts and templates",
    "Follow-up strategies"
  ]
}}

IMPORTANT INSTRUCTIONS:
1. If information is not available, use "Not specified" or leave arrays empty
2. Focus on extracting actionable information that would help AI sales agents
3. Return only valid JSON, no additional text

Return the JSON structure.
"""
        
        else:
            # Fallback to detection prompt for unknown types
            return self._build_detection_prompt(documents_text)
    
    def _build_detection_prompt(self, documents_text: str) -> str:
        """
        Build the original detection-based prompt (existing logic).
        """
        return f"""
You are an expert business analyst specializing in extracting structured knowledge from various types of documents. 
First, analyze the following documents to determine their type, then extract relevant information accordingly.

Documents to analyze:
{documents_text}

STEP 1: Determine the document type(s) from these categories:
- "company_info": Documents about a specific company (business plans, company profiles, product sheets)
- "sales_training": Training materials about sales techniques, methodologies, or best practices
- "industry_knowledge": Documents about industries, markets, or general business knowledge
- "product_docs": Technical documentation, user manuals, or product specifications
- "mixed": Documents containing multiple types of information

STEP 2: Extract knowledge based on the document type(s):

For COMPANY_INFO documents, extract:
{{
  "document_type": "company_info",
  "company_info": {{
    "company_name": "Name of the company",
    "industry": "Primary industry or sector",
    "company_size": "Size category (startup, SMB, enterprise)",
    "mission": "Company mission statement",
    "values": ["Core company values"],
    "founding_year": "Year founded if mentioned",
    "headquarters": "Location if mentioned"
  }},
  "sales_approach": "Detailed description of the company's sales methodology, approach, and philosophy",
  "products": [
    {{
      "name": "Product/service name",
      "description": "Detailed description",
      "target_customers": "Who this product serves",
      "key_features": ["Main features"],
      "benefits": ["Key benefits"],
      "pricing_model": "How it's priced if mentioned"
    }}
  ],
  "key_messages": [
    "Key messages the company uses in sales and marketing",
    "Value propositions",
    "Unique selling points"
  ],
  "value_propositions": [
    "Primary value propositions",
    "What makes them different",
    "Key benefits they provide"
  ],
  "target_audience": {{
    "primary_customers": "Description of ideal customers",
    "industries": ["Industries they serve"],
    "company_sizes": ["Target company sizes"],
    "pain_points": ["Problems they solve"]
  }},
  "competitive_advantages": [
    "What makes them unique",
    "Competitive differentiators",
    "Strengths vs competitors"
  ]
}}

For SALES_TRAINING documents, extract:
{{
  "document_type": "sales_training",
  "sales_methodologies": [
    "Sales techniques and approaches mentioned",
    "Best practices for prospecting",
    "Communication strategies"
  ],
  "target_audience_insights": {{
    "buyer_personas": ["Types of buyers mentioned"],
    "pain_points": ["Common problems prospects face"],
    "decision_making_process": ["How buyers make decisions"]
  }},
  "sales_philosophy": "Overall sales approach and mindset",
  "key_messages": [
    "Important sales messages and frameworks",
    "Value propositions that work",
    "Objection handling techniques"
  ],
  "industry_context": {{
    "industries_mentioned": ["Industries discussed"],
    "market_trends": ["Trends and insights"],
    "competitive_landscape": ["Competitive information"]
  }},
  "practical_advice": [
    "Actionable sales tips",
    "Scripts and templates",
    "Follow-up strategies"
  ]
}}

For INDUSTRY_KNOWLEDGE documents, extract:
{{
  "document_type": "industry_knowledge",
  "industry_insights": {{
    "industries_covered": ["Industries discussed"],
    "market_size": "Market information if mentioned",
    "growth_trends": ["Growth patterns and trends"],
    "key_players": ["Major companies mentioned"]
  }},
  "buyer_behavior": {{
    "decision_makers": ["Types of decision makers"],
    "purchase_process": ["How buying decisions are made"],
    "pain_points": ["Common industry challenges"]
  }},
  "sales_opportunities": [
    "Potential sales opportunities mentioned",
    "Timing considerations",
    "Market entry strategies"
  ],
  "competitive_intelligence": [
    "Competitive information",
    "Market positioning insights",
    "Differentiation strategies"
  ]
}}

For MIXED documents, extract the most relevant information from the appropriate categories above.

IMPORTANT INSTRUCTIONS:
1. If the document is clearly about sales training (like "how to be the best AI SDR"), use the sales_training format
2. If the document is about a specific company, use the company_info format
3. If the document is about industries or markets, use the industry_knowledge format
4. If information is not available, use "Not specified" or leave arrays empty
5. Focus on extracting actionable information that would help AI sales agents
6. Return only valid JSON, no additional text

Return the appropriate JSON structure based on the document type you identify.
"""
    
    def _parse_claude_response(self, response_text: str) -> Dict[str, Any]:
        """
        Parse Claude's response and extract the JSON.
        """
        try:
            # Try to find JSON in the response
            start_idx = response_text.find('{')
            end_idx = response_text.rfind('}') + 1
            
            if start_idx != -1 and end_idx != -1:
                json_str = response_text[start_idx:end_idx]
                return json.loads(json_str)
            else:
                logger.warning("No JSON found in Claude response")
                return self._get_default_knowledge_structure()
                
        except json.JSONDecodeError as e:
            logger.error(f"JSON parsing error: {e}")
            return self._get_default_knowledge_structure()
    
    def _get_default_knowledge_structure(self) -> Dict[str, Any]:
        """
        Return a default knowledge structure if extraction fails.
        """
        return {
            "document_type": "company_info",
            "company_info": {
                "company_name": "Not specified",
                "industry": "Not specified",
                "company_size": "Not specified",
                "mission": "Not specified",
                "values": [],
                "founding_year": "Not specified",
                "headquarters": "Not specified"
            },
            "sales_approach": "Not specified - please review and update",
            "products": [],
            "key_messages": [],
            "value_propositions": [],
            "target_audience": {
                "primary_customers": "Not specified",
                "industries": [],
                "company_sizes": [],
                "pain_points": []
            },
            "competitive_advantages": []
        }
    
    def _validate_and_score_knowledge(self, knowledge: Dict[str, Any], documents_content: List[Dict[str, Any]], document_type: str = None) -> Dict[str, Any]:
        """
        Validate extracted knowledge and add confidence scores.
        
        Args:
            knowledge: Extracted knowledge structure
            documents_content: Original document content for validation
            document_type: Document type for context
            
        Returns:
            Enhanced knowledge with quality metrics and confidence scores
        """
        logger.info("Validating and scoring extracted knowledge")
        
        # Calculate quality metrics
        quality_metrics = self._calculate_quality_metrics(knowledge, documents_content, document_type)
        
        # Add confidence scores to each section
        enhanced_knowledge = knowledge.copy()
        enhanced_knowledge["quality_metrics"] = quality_metrics
        
        # Add confidence scores to individual sections
        enhanced_knowledge = self._add_confidence_scores(enhanced_knowledge, quality_metrics)
        
        # Validate completeness
        completeness_score = self._assess_completeness(enhanced_knowledge, document_type)
        enhanced_knowledge["quality_metrics"]["completeness_score"] = completeness_score
        
        # Validate consistency
        consistency_score = self._assess_consistency(enhanced_knowledge)
        enhanced_knowledge["quality_metrics"]["consistency_score"] = consistency_score
        
        # Overall quality score
        overall_score = self._calculate_overall_quality_score(quality_metrics, completeness_score, consistency_score)
        enhanced_knowledge["quality_metrics"]["overall_score"] = overall_score
        
        logger.info(f"Knowledge validation completed - Overall score: {overall_score:.2f}")
        
        return enhanced_knowledge
    
    def _calculate_quality_metrics(self, knowledge: Dict[str, Any], documents_content: List[Dict[str, Any]], document_type: str = None) -> Dict[str, Any]:
        """Calculate various quality metrics for extracted knowledge"""
        metrics = {}
        
        # Document coverage
        total_content_length = sum(len(doc.get('content', '')) for doc in documents_content)
        metrics["document_coverage"] = {
            "total_chars": total_content_length,
            "documents_processed": len(documents_content),
            "chunks_processed": sum(1 for doc in documents_content if doc.get('is_chunk', False))
        }
        
        # Content richness
        metrics["content_richness"] = self._assess_content_richness(knowledge)
        
        # Specificity analysis
        metrics["specificity"] = self._assess_specificity(knowledge)
        
        # Actionability score
        metrics["actionability"] = self._assess_actionability(knowledge)
        
        # Document type alignment
        if document_type:
            metrics["type_alignment"] = self._assess_type_alignment(knowledge, document_type)
        
        return metrics
    
    def _assess_content_richness(self, knowledge: Dict[str, Any]) -> Dict[str, float]:
        """Assess how rich and detailed the extracted content is"""
        richness = {}
        
        # Company info richness
        company_info = knowledge.get("company_info", {})
        richness["company_info"] = len([v for v in company_info.values() if v and v != "Not specified"]) / max(len(company_info), 1)
        
        # Products richness
        products = knowledge.get("products", [])
        richness["products"] = min(len(products) / 3, 1.0)  # Normalize to max 3 products
        
        # Value propositions richness
        value_props = knowledge.get("value_propositions", [])
        richness["value_propositions"] = min(len(value_props) / 5, 1.0)  # Normalize to max 5
        
        # Target audience richness
        target_audience = knowledge.get("target_audience", {})
        if isinstance(target_audience, dict):
            total_audience_items = sum(len(v) if isinstance(v, list) else 1 for v in target_audience.values())
            richness["target_audience"] = min(total_audience_items / 10, 1.0)  # Normalize to max 10 items
        else:
            richness["target_audience"] = 0.0
        
        # Overall richness
        richness["overall"] = sum(richness.values()) / len(richness)
        
        return richness
    
    def _assess_specificity(self, knowledge: Dict[str, Any]) -> Dict[str, float]:
        """Assess how specific vs generic the extracted information is"""
        specificity = {}
        
        # Check for generic terms
        generic_terms = ["not specified", "unknown", "various", "multiple", "several", "some"]
        
        # Company info specificity
        company_info = knowledge.get("company_info", {})
        specific_fields = 0
        total_fields = len(company_info)
        
        for value in company_info.values():
            if isinstance(value, str) and not any(term in value.lower() for term in generic_terms):
                specific_fields += 1
        
        specificity["company_info"] = specific_fields / max(total_fields, 1)
        
        # Sales approach specificity
        sales_approach = knowledge.get("sales_approach", "")
        if sales_approach and not any(term in sales_approach.lower() for term in generic_terms):
            specificity["sales_approach"] = 1.0
        else:
            specificity["sales_approach"] = 0.0
        
        # Products specificity
        products = knowledge.get("products", [])
        specific_products = 0
        for product in products:
            if isinstance(product, str) and not any(term in product.lower() for term in generic_terms):
                specific_products += 1
        
        specificity["products"] = specific_products / max(len(products), 1)
        
        # Overall specificity
        specificity["overall"] = sum(specificity.values()) / len(specificity)
        
        return specificity
    
    def _assess_actionability(self, knowledge: Dict[str, Any]) -> Dict[str, float]:
        """Assess how actionable the extracted knowledge is for sales activities"""
        actionability = {}
        
        # Value propositions actionability
        value_props = knowledge.get("value_propositions", [])
        actionable_vps = 0
        for vp in value_props:
            if isinstance(vp, str):
                # Check for actionable language (ROI, benefits, outcomes)
                actionable_keywords = ["increase", "reduce", "save", "improve", "boost", "enhance", "streamline", "automate"]
                if any(keyword in vp.lower() for keyword in actionable_keywords):
                    actionable_vps += 1
        
        actionability["value_propositions"] = actionable_vps / max(len(value_props), 1)
        
        # Target audience actionability
        target_audience = knowledge.get("target_audience", {})
        if isinstance(target_audience, dict):
            roles = target_audience.get("roles", [])
            industries = target_audience.get("industries", [])
            
            # More specific roles and industries are more actionable
            actionability["target_audience"] = min((len(roles) + len(industries)) / 6, 1.0)
        else:
            actionability["target_audience"] = 0.0
        
        # Competitive advantages actionability
        competitive_advantages = knowledge.get("competitive_advantages", [])
        actionable_advantages = 0
        for advantage in competitive_advantages:
            if isinstance(advantage, str):
                # Check for specific differentiators
                if len(advantage.split()) > 3:  # More than 3 words suggests specificity
                    actionable_advantages += 1
        
        actionability["competitive_advantages"] = actionable_advantages / max(len(competitive_advantages), 1)
        
        # Overall actionability
        actionability["overall"] = sum(actionability.values()) / len(actionability)
        
        return actionability
    
    def _assess_type_alignment(self, knowledge: Dict[str, Any], document_type: str) -> float:
        """Assess how well the extracted knowledge aligns with the document type"""
        alignment_score = 0.0
        
        if document_type == "sales_training":
            # Check for sales-specific content
            sales_indicators = ["sales", "selling", "prospect", "customer", "close", "deal", "pipeline"]
            content_text = str(knowledge).lower()
            alignment_score = sum(1 for indicator in sales_indicators if indicator in content_text) / len(sales_indicators)
        
        elif document_type == "company_info":
            # Check for company-specific content
            company_indicators = ["company", "business", "organization", "founded", "mission", "vision", "values"]
            content_text = str(knowledge).lower()
            alignment_score = sum(1 for indicator in company_indicators if indicator in content_text) / len(company_indicators)
        
        elif document_type == "industry_knowledge":
            # Check for industry-specific content
            industry_indicators = ["industry", "market", "trend", "competitor", "regulation", "standard"]
            content_text = str(knowledge).lower()
            alignment_score = sum(1 for indicator in industry_indicators if indicator in content_text) / len(industry_indicators)
        
        return min(alignment_score, 1.0)
    
    def _add_confidence_scores(self, knowledge: Dict[str, Any], quality_metrics: Dict[str, Any]) -> Dict[str, Any]:
        """Add confidence scores to each section of the knowledge"""
        enhanced = knowledge.copy()
        
        # Base confidence on quality metrics
        richness = quality_metrics.get("content_richness", {})
        specificity = quality_metrics.get("specificity", {})
        actionability = quality_metrics.get("actionability", {})
        
        # Calculate confidence for each section
        sections = ["company_info", "products", "value_propositions", "target_audience", "competitive_advantages"]
        
        for section in sections:
            if section in enhanced:
                # Weighted confidence based on richness, specificity, and actionability
                section_richness = richness.get(section, 0.5)
                section_specificity = specificity.get(section, 0.5)
                section_actionability = actionability.get(section, 0.5)
                
                confidence = (section_richness * 0.4 + section_specificity * 0.4 + section_actionability * 0.2)
                
                if isinstance(enhanced[section], dict):
                    enhanced[section]["confidence"] = confidence
                else:
                    enhanced[section] = {
                        "content": enhanced[section],
                        "confidence": confidence
                    }
        
        return enhanced
    
    def _assess_completeness(self, knowledge: Dict[str, Any], document_type: str = None) -> float:
        """Assess how complete the extracted knowledge is"""
        required_fields = {
            "company_info": ["company_name", "industry"],
            "products": [],
            "value_propositions": [],
            "target_audience": ["roles"],
            "sales_approach": [],
            "competitive_advantages": []
        }
        
        # Adjust requirements based on document type
        if document_type == "sales_training":
            required_fields["sales_approach"] = ["methodology"]
            required_fields["value_propositions"] = ["primary_benefits"]
        elif document_type == "company_info":
            required_fields["company_info"] = ["company_name", "industry", "company_size"]
        
        total_required = 0
        completed_fields = 0
        
        for section, fields in required_fields.items():
            if section in knowledge:
                if fields:  # Has specific required fields
                    for field in fields:
                        total_required += 1
                        if field in knowledge[section] and knowledge[section][field] and knowledge[section][field] != "Not specified":
                            completed_fields += 1
                else:  # Just needs to exist and have content
                    total_required += 1
                    if knowledge[section] and knowledge[section] != "Not specified":
                        completed_fields += 1
        
        return completed_fields / max(total_required, 1)
    
    def _assess_consistency(self, knowledge: Dict[str, Any]) -> float:
        """Assess internal consistency of the extracted knowledge"""
        consistency_score = 1.0  # Start with perfect consistency
        
        # Check for contradictions in company info
        company_info = knowledge.get("company_info", {})
        if isinstance(company_info, dict):
            # Check for conflicting company sizes
            size_indicators = ["startup", "small", "medium", "large", "enterprise"]
            size_values = [str(v).lower() for v in company_info.values()]
            
            conflicting_sizes = []
            for size in size_indicators:
                if sum(1 for val in size_values if size in val) > 1:
                    conflicting_sizes.append(size)
            
            if conflicting_sizes:
                consistency_score -= 0.2
        
        # Check for conflicting industry information
        industry_info = str(company_info).lower()
        tech_terms = ["tech", "software", "saas", "ai", "technology"]
        non_tech_terms = ["healthcare", "finance", "retail", "manufacturing"]
        
        has_tech = any(term in industry_info for term in tech_terms)
        has_non_tech = any(term in industry_info for term in non_tech_terms)
        
        if has_tech and has_non_tech:
            consistency_score -= 0.1
        
        return max(consistency_score, 0.0)
    
    def _calculate_overall_quality_score(self, quality_metrics: Dict[str, Any], completeness_score: float, consistency_score: float) -> float:
        """Calculate overall quality score from all metrics"""
        # Weight different aspects
        richness_score = quality_metrics.get("content_richness", {}).get("overall", 0.5)
        specificity_score = quality_metrics.get("specificity", {}).get("overall", 0.5)
        actionability_score = quality_metrics.get("actionability", {}).get("overall", 0.5)
        type_alignment_score = quality_metrics.get("type_alignment", 0.5)
        
        # Weighted average
        overall_score = (
            richness_score * 0.25 +
            specificity_score * 0.25 +
            actionability_score * 0.25 +
            completeness_score * 0.15 +
            consistency_score * 0.10
        )
        
        return min(overall_score, 1.0)
    
    def _chunk_large_document(self, content: str, file_path: str, max_chunk_size: int = 100000) -> List[Dict[str, Any]]:
        """
        Split large documents into manageable chunks for processing.
        
        Args:
            content: Document content to chunk
            file_path: Original file path for metadata
            max_chunk_size: Maximum characters per chunk
            
        Returns:
            List of document chunks with metadata
        """
        logger.info(f"Chunking document {Path(file_path).name} ({len(content)} chars)")
        
        chunks = []
        
        # Try to split by logical sections first
        logical_chunks = self._split_by_logical_sections(content, max_chunk_size)
        
        if len(logical_chunks) > 1:
            logger.info(f"Split into {len(logical_chunks)} logical sections")
            for i, chunk_content in enumerate(logical_chunks):
                chunks.append({
                    'filename': Path(file_path).name,
                    'content': chunk_content,
                    'is_chunk': True,
                    'chunk_index': i,
                    'total_chunks': len(logical_chunks),
                    'chunk_type': 'logical_section'
                })
        else:
            # Fall back to paragraph-based chunking
            paragraph_chunks = self._split_by_paragraphs(content, max_chunk_size)
            logger.info(f"Split into {len(paragraph_chunks)} paragraph-based chunks")
            for i, chunk_content in enumerate(paragraph_chunks):
                chunks.append({
                    'filename': Path(file_path).name,
                    'content': chunk_content,
                    'is_chunk': True,
                    'chunk_index': i,
                    'total_chunks': len(paragraph_chunks),
                    'chunk_type': 'paragraph'
                })
        
        return chunks
    
    def _split_by_logical_sections(self, content: str, max_chunk_size: int) -> List[str]:
        """Split content by logical sections (headers, chapters, etc.)"""
        chunks = []
        
        # Common section markers
        section_patterns = [
            r'\n(?:Chapter|Section|Part)\s+\d+',
            r'\n#{1,6}\s+',  # Markdown headers
            r'\n[A-Z][A-Z\s]{10,}\n',  # ALL CAPS headers
            r'\n\d+\.\s+[A-Z]',  # Numbered sections
            r'\n(?:Introduction|Overview|Summary|Conclusion)',
        ]
        
        # Find section boundaries
        boundaries = []
        for pattern in section_patterns:
            matches = list(re.finditer(pattern, content, re.IGNORECASE))
            boundaries.extend([(m.start(), m.group()) for m in matches])
        
        # Sort boundaries by position
        boundaries.sort(key=lambda x: x[0])
        
        if not boundaries:
            return [content]  # No logical sections found
        
        # Create chunks based on boundaries
        start = 0
        for i, (pos, marker) in enumerate(boundaries):
            if pos - start > max_chunk_size:
                # Current chunk is too large, split it
                chunk_content = content[start:pos]
                chunks.append(chunk_content)
                start = pos
        
        # Add remaining content
        if start < len(content):
            chunks.append(content[start:])
        
        return chunks if chunks else [content]
    
    def _split_by_paragraphs(self, content: str, max_chunk_size: int) -> List[str]:
        """Split content by paragraphs when logical sections aren't available"""
        paragraphs = content.split('\n\n')
        chunks = []
        current_chunk = ""
        
        for paragraph in paragraphs:
            # If adding this paragraph would exceed max size, start new chunk
            if len(current_chunk) + len(paragraph) > max_chunk_size and current_chunk:
                chunks.append(current_chunk.strip())
                current_chunk = paragraph
            else:
                current_chunk += "\n\n" + paragraph if current_chunk else paragraph
        
        # Add final chunk
        if current_chunk:
            chunks.append(current_chunk.strip())
        
        return chunks if chunks else [content]
    
    async def extract_knowledge_parallel(self, file_paths: List[str], document_type: str = None) -> Dict[str, Any]:
        """
        Process multiple documents in parallel for better performance.
        
        Args:
            file_paths: List of file paths to process
            document_type: Optional document type to guide extraction
            
        Returns:
            Dict with aggregated knowledge from all documents
        """
        logger.info(f"Starting parallel processing of {len(file_paths)} files")
        
        try:
            # Use ThreadPoolExecutor for I/O operations
            with ThreadPoolExecutor(max_workers=min(len(file_paths), 4)) as executor:
                # Process each file
                futures = []
                for file_path in file_paths:
                    future = executor.submit(self._process_single_file, file_path, document_type)
                    futures.append(future)
                
                # Collect results
                results = []
                for future in futures:
                    try:
                        result = future.result(timeout=300)  # 5 minute timeout per file
                        if result and result.get('success'):
                            results.append(result['knowledge'])
                    except Exception as e:
                        logger.error(f"Error processing file in parallel: {e}")
                        continue
            
            if not results:
                return {
                    "success": False,
                    "error": "No documents processed successfully"
                }
            
            # Aggregate knowledge from all documents
            aggregated_knowledge = self._aggregate_knowledge_from_chunks(results)
            
            logger.info(f"Parallel processing completed: {len(results)} documents processed")
            return {
                "success": True,
                "knowledge": aggregated_knowledge,
                "processing_method": "parallel",
                "documents_processed": len(results)
            }
            
        except Exception as e:
            logger.error(f"Parallel processing error: {e}")
            return {
                "success": False,
                "error": str(e)
            }
    
    def _process_single_file(self, file_path: str, document_type: str = None) -> Dict[str, Any]:
        """Process a single file (used in parallel processing)"""
        try:
            content = self._read_document(file_path)
            if not content:
                return {"success": False, "error": "Could not read file"}
            
            # Check if chunking is needed
            if len(content) > 100000:
                chunks = self._chunk_large_document(content, file_path)
                # Process chunks and aggregate
                chunk_results = []
                for chunk in chunks:
                    chunk_knowledge = self._extract_with_claude([chunk], document_type)
                    chunk_results.append(chunk_knowledge)
                
                aggregated = self._aggregate_knowledge_from_chunks(chunk_results)
                return {"success": True, "knowledge": aggregated}
            else:
                # Process normally
                documents_content = [{
                    'filename': Path(file_path).name,
                    'content': content,
                    'is_chunk': False,
                    'chunk_index': 0
                }]
                
                knowledge = self._extract_with_claude(documents_content, document_type)
                return {"success": True, "knowledge": knowledge}
                
        except Exception as e:
            logger.error(f"Error processing single file {file_path}: {e}")
            return {"success": False, "error": str(e)}
    
    def _aggregate_knowledge_from_chunks(self, knowledge_chunks: List[Dict[str, Any]]) -> Dict[str, Any]:
        """
        Aggregate knowledge from multiple chunks or documents.
        
        Args:
            knowledge_chunks: List of knowledge dictionaries from chunks
            
        Returns:
            Aggregated knowledge structure
        """
        if not knowledge_chunks:
            return self._get_default_knowledge_structure()
        
        if len(knowledge_chunks) == 1:
            return knowledge_chunks[0]
        
        logger.info(f"Aggregating knowledge from {len(knowledge_chunks)} chunks")
        
        # Initialize aggregated structure
        aggregated = {
            "company_info": {},
            "sales_approach": "",
            "products": [],
            "key_messages": [],
            "value_propositions": [],
            "target_audience": {
                "roles": [],
                "industries": [],
                "company_sizes": [],
                "pain_points": []
            },
            "competitive_advantages": []
        }
        
        # Aggregate each section
        for chunk in knowledge_chunks:
            # Company info - merge dictionaries
            if chunk.get("company_info"):
                for key, value in chunk["company_info"].items():
                    if value and value != "Not specified":
                        aggregated["company_info"][key] = value
            
            # Sales approach - combine approaches
            if chunk.get("sales_approach") and chunk["sales_approach"] != "Not specified":
                if aggregated["sales_approach"]:
                    aggregated["sales_approach"] += f"; {chunk['sales_approach']}"
                else:
                    aggregated["sales_approach"] = chunk["sales_approach"]
            
            # Products - extend lists
            if chunk.get("products"):
                aggregated["products"].extend(chunk["products"])
            
            # Key messages - extend lists
            if chunk.get("key_messages"):
                aggregated["key_messages"].extend(chunk["key_messages"])
            
            # Value propositions - extend lists
            if chunk.get("value_propositions"):
                aggregated["value_propositions"].extend(chunk["value_propositions"])
            
            # Target audience - merge lists
            if chunk.get("target_audience"):
                for key, value in chunk["target_audience"].items():
                    if isinstance(value, list):
                        aggregated["target_audience"][key].extend(value)
                    elif isinstance(value, dict):
                        aggregated["target_audience"][key].update(value)
            
            # Competitive advantages - extend lists
            if chunk.get("competitive_advantages"):
                aggregated["competitive_advantages"].extend(chunk["competitive_advantages"])
        
        # Remove duplicates and clean up
        aggregated = self._clean_aggregated_knowledge(aggregated)
        
        return aggregated
    
    def _clean_aggregated_knowledge(self, knowledge: Dict[str, Any]) -> Dict[str, Any]:
        """Clean and deduplicate aggregated knowledge"""
        # Remove duplicates from lists
        for key in ["products", "key_messages", "value_propositions", "competitive_advantages"]:
            if key in knowledge and isinstance(knowledge[key], list):
                knowledge[key] = list(set(knowledge[key]))
        
        # Clean target audience lists
        if "target_audience" in knowledge:
            for sub_key in ["roles", "industries", "company_sizes", "pain_points"]:
                if sub_key in knowledge["target_audience"] and isinstance(knowledge["target_audience"][sub_key], list):
                    knowledge["target_audience"][sub_key] = list(set(knowledge["target_audience"][sub_key]))
        
        return knowledge
    
    def get_agent_info(self) -> Dict[str, Any]:
        """
        Get information about the Knowledge Extraction Agent.
        """
        return {
            "name": "Knowledge Extraction Agent",
            "description": "Uses Claude AI to extract structured knowledge from company documents with advanced chunking and parallel processing",
            "capabilities": [
                "Document analysis and processing",
                "Company information extraction",
                "Sales approach identification",
                "Product knowledge extraction",
                "Value proposition identification",
                "Target audience analysis",
                "Large document chunking (>100k tokens)",
                "Parallel document processing",
                "Intelligent knowledge aggregation",
                "Logical section detection"
            ],
            "supported_formats": ["PDF", "Word", "PowerPoint", "Text"],
            "ai_model": "Claude-3-Sonnet",
            "specialization": "Business document analysis and knowledge structuring with scalability",
            "processing_methods": ["standard", "chunked", "parallel"],
            "max_document_size": "Unlimited (with chunking)",
            "parallel_workers": 4
        }
